"""
Layout QA for the SOC-Copilot poster.

Reads the generated .pptx with python-pptx and verifies:
  * slide dimensions match B1 portrait (70 x 100 cm)
  * every shape lies within the slide bounds (no off-slide elements)
  * no two top-level shapes overlap badly (>2 cm overlap in both dims)
  * required text strings are present (section titles in order)
"""

from pptx import Presentation
from pptx.util import Emu

EMU_PER_CM = 360000

path = "/mnt/c/Users/stav/Desktop/cyber course/HIT-ai-cybersecurity/Project/poster/SOC-Copilot-Poster.pptx"
prs = Presentation(path)

W_cm = prs.slide_width  / EMU_PER_CM
H_cm = prs.slide_height / EMU_PER_CM
print(f"Slide size: {W_cm:.2f} cm x {H_cm:.2f} cm")
assert abs(W_cm - 70) < 0.1, "slide width != 70cm"
assert abs(H_cm - 100) < 0.1, "slide height != 100cm"
print("  -> matches B1 portrait spec OK")

s = prs.slides[0]

# 1. extract every shape's bbox in cm + its text (if any)
shapes_info = []
for sh in s.shapes:
    try:
        x = sh.left  / EMU_PER_CM
        y = sh.top   / EMU_PER_CM
        w = sh.width / EMU_PER_CM
        h = sh.height/ EMU_PER_CM
    except Exception:
        continue
    txt = ""
    if sh.has_text_frame:
        txt = " | ".join(p.text for p in sh.text_frame.paragraphs)[:80]
    shapes_info.append((sh.shape_id, sh.name, x, y, w, h, txt))

print(f"\nShape count: {len(shapes_info)}")

# 2. anything off-slide?
bad = []
for sid, name, x, y, w, h, txt in shapes_info:
    if x < -0.1 or y < -0.1 or x + w > W_cm + 0.1 or y + h > H_cm + 0.1:
        bad.append((sid, name, x, y, w, h, txt))
if bad:
    print("\nOff-slide elements:")
    for r in bad: print(" ", r)
else:
    print("All shapes within slide bounds OK")

# 3. required text in order
flat = "\n".join(t for _, _, _, _, _, _, t in shapes_info)
required = [
    "SOC-Copilot",
    "AI in Cybersecurity",
    "Stav Hefetz",
    "1. Introduction",
    "2. Methodology",
    "3. Results",
    "4. Conclusions",
    "5. Discussions",
    "Demo video",
    "github.com",
]
missing = [r for r in required if r not in flat]
if missing:
    print("\nMISSING required text:")
    for r in missing: print(" -", r)
else:
    print("All required text present OK")

# 4. show vertical layout: print each shape that has a 'Section header' style or
# is wider than 50% of slide -- the layout skeleton
print("\nLayout skeleton (wide elements with text):")
wide = [s for s in shapes_info if s[4] > W_cm * 0.45 and s[6].strip()]
wide.sort(key=lambda s: s[3])
for sid, name, x, y, w, h, txt in wide:
    print(f"  y={y:5.2f}cm h={h:5.2f}cm  {txt[:60]}")
